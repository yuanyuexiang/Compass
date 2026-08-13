import io

import fitz
from pptx import Presentation
from pptx.util import Inches

from app.parsing.documents import parse_attachment, pdf_to_text, pptx_to_text


def make_pdf(text: str = "") -> bytes:
    doc = fitz.open()
    page = doc.new_page()
    if text:
        page.insert_text((72, 72), text)
    return doc.tobytes()


def test_pdf_with_text_layer():
    data = make_pdf("Purchase budget: 3,300,000 CNY. Deadline 2026-07-09.")
    text, needs_ocr = pdf_to_text(data)
    assert "3,300,000" in text
    assert "[第1页]" in text
    assert needs_ocr is False


def test_scanned_pdf_flagged_for_ocr():
    text, needs_ocr = pdf_to_text(make_pdf())  # 无文本层 → 疑似扫描件
    assert needs_ocr is True


def test_parse_attachment_dispatch():
    body = "Tender notice content. " * 5  # 需超过扫描件判定阈值（30 字符/页）
    text, needs_ocr = parse_attachment("公告.pdf", make_pdf(body))
    assert "Tender notice" in text and needs_ocr is False
    text, needs_ocr = parse_attachment("其他.zip", b"")
    assert text == "" and needs_ocr is False


def test_pptx_keeps_slide_numbers_and_tables():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "智慧园区解决方案"
    table = slide.shapes.add_table(1, 2, 100, 100, 500, 100).table
    table.cell(0, 0).text = "能力"
    table.cell(0, 1).text = "AI 安防"
    group = slide.shapes.add_group_shape()
    grouped_box = group.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1))
    grouped_box.text = "成组文本：数据治理"
    stream = io.BytesIO()
    presentation.save(stream)

    text = pptx_to_text(stream.getvalue())
    assert "[第1页]" in text
    assert "智慧园区解决方案" in text
    assert "能力 | AI 安防" in text
    assert "成组文本：数据治理" in text

    dispatched, needs_ocr = parse_attachment("企业介绍.PPTX", stream.getvalue())
    assert dispatched == text
    assert needs_ocr is False
