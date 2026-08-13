"""附件解析：PDF/Word → 带页码标记的纯文本，扫描件标记 OCR（tech-design.md §4.2）。"""

import io

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

# 平均每页文本字符数低于此阈值判定为扫描件
SCANNED_CHARS_PER_PAGE = 30


def pdf_to_text(data: bytes) -> tuple[str, bool]:
    """返回 (带页码标记的文本, 是否疑似扫描件需 OCR)。"""
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        pages = [page.get_text().strip() for page in doc]
        total_chars = sum(len(p) for p in pages)
        needs_ocr = len(pages) > 0 and total_chars / len(pages) < SCANNED_CHARS_PER_PAGE
        text = "\n\n".join(f"[第{i}页]\n{p}" for i, p in enumerate(pages, 1) if p)
        return text, needs_ocr
    finally:
        doc.close()


def docx_to_text(data: bytes) -> str:
    doc = Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def pptx_to_text(data: bytes) -> str:
    """提取幻灯片正文、表格与演讲者备注，并保留可追溯的页码标记。"""
    presentation = Presentation(io.BytesIO(data))
    pages: list[str] = []
    for page_number, slide in enumerate(presentation.slides, 1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
        if slide.has_notes_slide:
            notes = [
                shape.text.strip()
                for shape in slide.notes_slide.notes_text_frame.paragraphs
                if shape.text.strip()
            ]
            if notes:
                parts.append("演讲者备注：\n" + "\n".join(notes))
        if parts:
            pages.append(f"[第{page_number}页]\n" + "\n".join(parts))
    return "\n\n".join(pages)


def ocr_pdf(data: bytes) -> str:
    """扫描件 OCR 兜底（PaddleOCR，PP-OCRv6）。

    M1 后期接入：PaddleOCR 依赖较重（paddlepaddle 需官方源），单独装在 worker 镜像。
    当前先抛错并由流水线将附件标记为 needs_ocr，不阻塞主链路。
    """
    raise NotImplementedError("OCR 尚未接入（见 tech-design.md 附录 D 遗留项）")


def parse_attachment(filename: str, data: bytes) -> tuple[str, bool]:
    """按扩展名分发解析器，返回 (文本, needs_ocr)。不支持的类型返回空文本。"""
    lower = filename.lower()
    if lower.endswith(".pdf"):
        return pdf_to_text(data)
    if lower.endswith(".docx"):
        return docx_to_text(data), False
    if lower.endswith(".pptx"):
        return pptx_to_text(data), False
    if lower.endswith(".txt"):
        return data.decode("utf-8", errors="replace"), False
    # .doc 老格式需 LibreOffice 转换（M1 遗留项）；其余类型暂不解析
    return "", False
